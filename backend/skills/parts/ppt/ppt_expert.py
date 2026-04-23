import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt

class PPTMasterSkill:
    """
    Industrial PPT Design Master: High-fidelity layout and content orchestration.
    Powered by python-pptx for the Nexus Center platform.
    """
    def __init__(self, ppt_path: str = None):
        self.prs = Presentation(ppt_path) if ppt_path else Presentation()

    def add_elegant_slide(self, title_text: str, body_text: str):
        # Using a default title/content layout for the Pro Max version
        slide_layout = self.prs.slide_layouts[1] 
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = title_text
        content.text = body_text
        
        # Apply Industrial Zenith styling
        for paragraph in content.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(20)
                run.font.name = 'Outfit' # Matches our UI aesthetic

    def save_presentation(self, output_path: str):
        self.prs.save(output_path)
        return {"status": "success", "file": output_path}

if __name__ == "__main__":
    try:
        # Standard ACP JSON Over-STDIN Pipe
        input_data = json.load(sys.stdin)
        action = input_data.get("action")
        
        master = PPTMasterSkill(input_data.get("path"))
        
        if action == "add_slide":
            master.add_elegant_slide(input_data["title"], input_data["body"])
            res = master.save_presentation(input_data["output"])
            print(json.dumps(res))
        elif action == "analyze_layout":
            # Simulation of design-intent analysis
            print(json.dumps({"status": "ready", "slides_count": len(master.prs.slides)}))
            
    except Exception as e:
        print(json.dumps({"error": str(e)}), file=sys.stderr)
        sys.exit(1)
