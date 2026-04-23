import sys
import json
import pandas as pd
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook

class OmniBridgeSkill:
    """
    Industrial Omni-Bridge: Orchestrating Data & Visuals across the Office Suite.
    Powered by the Nexus Double-Engine Paradigm.
    """
    def __init__(self):
        pass

    def excel_to_ppt_sync(self, excel_path, ppt_path, sheet_name, table_range):
        """
        Extracts data from Excel and updates/inserts into a PPT slide.
        """
        # 1. Read Data (Precision Logic)
        df = pd.read_excel(excel_path, sheet_name=sheet_name)
        data_subset = df.to_records(index=False).tolist()
        
        # 2. Open PPT (Master Vision)
        prs = Presentation(ppt_path)
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
        slide.shapes.title.text = f"Data Sync from {sheet_name}"
        
        # 3. Create Native Table in PPT
        rows, cols = len(data_subset) + 1, len(df.columns)
        table = slide.shapes.add_table(rows, cols, left=0, top=1000000, width=6000000, height=3000000).table
        
        # Insert Headers
        for i, col_name in enumerate(df.columns):
            table.cell(0, i).text = str(col_name)
            
        # Insert Data
        for row_idx, row_data in enumerate(data_subset, 1):
            for col_idx, value in enumerate(row_data):
                table.cell(row_idx, col_idx).text = str(value)
        
        prs.save(ppt_path)
        return {"status": "success", "synced_rows": len(data_subset)}

if __name__ == "__main__":
    try:
        # Standard ACP Data Flow
        input_data = json.load(sys.stdin)
        bridge = OmniBridgeSkill()
        
        res = bridge.excel_to_ppt_sync(
            input_data["excel_path"],
            input_data["ppt_path"],
            input_data["sheet_name"],
            input_data.get("range")
        )
        print(json.dumps(res))
    except Exception as e:
        print(json.dumps({"error": str(e)}), file=sys.stderr)
        sys.exit(1)
