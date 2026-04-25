"""
Omni-Office Zenith: Industrial-Grade Skill Integration Tests.
Focuses on high-fidelity, end-to-end professional workflows (BDD style) 
rather than trivial HTTP assertions. High quality over high quantity.
"""
import sys
import os
import openpyxl
import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches

# Ensure src/skills and its shared folder are on the path
base_dir = os.path.join(os.path.dirname(__file__), '..')
sdk_dir = os.path.join(base_dir, '..', 'sdk')
sys.path.insert(0, base_dir)
sys.path.insert(0, os.path.join(base_dir, 'shared'))
sys.path.insert(0, os.path.join(sdk_dir, 'core', 'transport', 'acp'))

@pytest.fixture
def client():
    """Return a TestClient for the FastAPI skill bridge."""
    try:
        from httpx import Client  # type: ignore
        from skill_bridge import app  # type: ignore
        from starlette.testclient import TestClient  # type: ignore
        return TestClient(app)
    except ImportError:
        pytest.skip("skill_bridge not yet available")

class TestNexusCoreEngines:
    """Validates the core semantic and transport engines."""
    
    def test_vector_nexus_initialization(self):
        """Ensures the semantic search engine boots correctly even in offline mode."""
        from vector_nexus import VectorNexus  # type: ignore
        assert VectorNexus(api_key='') is not None

    def test_bridge_health_and_capabilities(self, client):
        """Verifies the Bridge API registers all expected Omni-Office experts."""
        response = client.get("/health")
        assert response.status_code == 200
        skills = response.json().get("skills", [])
        assert all(s in skills for s in ["excel", "ppt", "word", "vector-search"])


class TestFinancialModelingWorkflow:
    """Real-world scenarios for the Excel Expert."""
    
    def test_end_to_end_financial_report_and_pivot(self, client, tmp_path):
        """
        Scenario: An agent generates a quarterly financial report, 
        injects formulas, applies corporate formatting, and builds a Pivot Table.
        """
        output_path = tmp_path / "Q3_Financial_Report.xlsx"
        
        payload = {
            "output_path": str(output_path),
            "changes": [
                {"op": "set_value", "cell": "Sheet!A1", "value": "Region"},
                {"op": "set_value", "cell": "Sheet!B1", "value": "Revenue"},
                {"op": "set_value", "cell": "Sheet!A2", "value": "North America"},
                {"op": "set_value", "cell": "Sheet!B2", "value": 500000},
                {"op": "set_value", "cell": "Sheet!A3", "value": "EMEA"},
                {"op": "set_value", "cell": "Sheet!B3", "value": 300000},
                # Formula injection
                {"op": "set_formula", "cell": "Sheet!B4", "formula": "=SUM(B2:B3)"},
                # Formatting
                {"op": "format_range", "range": "Sheet!A1:B1", "bold": True, "fill_color": "008CA1"},
                # Pivot Summary
                {
                    "op": "create_pivottable",
                    "source": "Sheet!A1:B3",
                    "destination": "Summary!A1",
                    "rows": ["Region"],
                    "values": [{"field": "Revenue", "func": "SUM"}]
                }
            ],
            "office_context": {"activeSheet": "Sheet"}
        }

        response = client.post("/skills/excel", json=payload)
        assert response.status_code == 200, f"Bridge failed: {response.text}"
        
        # Verify Business Logic
        wb = openpyxl.load_workbook(output_path, data_only=False)
        assert "Summary" in wb.sheetnames, "Pivot table sheet was not created"
        assert wb["Sheet"]["B4"].value == "=SUM(B2:B3)", "Formula was not injected correctly"
        
        # Verify Pivot Aggregation (North America + EMEA)
        # B2 should contain the SUM of North America (500000)
        assert wb["Summary"]["B2"].value == 500000.0, "Pivot table aggregation failed"

    def test_legacy_format_and_error_handling(self, client, tmp_path):
        """
        Scenario: An agent attempts to output to an unsupported format (e.g., .txt)
        or reference a non-existent sheet in a formula. The system must fail fast and safely.
        """
        payload = {
            "output_path": str(tmp_path / "bad-formula.xlsx"),
            "changes": [{"op": "set_formula", "cell": "A1", "formula": "=GhostSheet!B2"}],
        }
        response = client.post("/skills/excel", json=payload)
        assert response.status_code == 500
        assert "Excel validation failed" in response.json()["detail"]

    def test_automated_invoice_generation_with_merged_headers_and_dynamic_widths(self, client, tmp_path):
        """
        Scenario: An agent builds a formatted invoice. This tests advanced
        operations like merge_cells, column width adjustment, and multi-cell formatting.
        """
        output_path = tmp_path / "Invoice.xlsx"
        
        payload = {
            "output_path": str(output_path),
            "changes": [
                {"op": "set_value", "cell": "Sheet!A1", "value": "CORPORATE INVOICE"},
                {"op": "merge_cells", "range": "Sheet!A1:D1"},
                {"op": "format_range", "range": "Sheet!A1:D1", "bold": True, "fill_color": "1E2761"},
                {"op": "set_column_width", "column": "A", "width": 20.0},
                {"op": "set_column_width", "column": "B", "width": 15.0},
                {"op": "add_header_row", "sheet": "Sheet", "row": 3, "headers": ["Item", "Qty", "Unit Price", "Total"]},
                {"op": "set_value", "cell": "Sheet!A4", "value": "Consulting"},
                {"op": "set_value", "cell": "Sheet!B4", "value": 10},
                {"op": "set_value", "cell": "Sheet!C4", "value": 150},
                {"op": "set_formula", "cell": "Sheet!D4", "formula": "=B4*C4"}
            ]
        }

        response = client.post("/skills/excel", json=payload)
        assert response.status_code == 200, f"Bridge failed: {response.text}"
        
        wb = openpyxl.load_workbook(output_path)
        ws = wb["Sheet"]
        # Check merge
        merged = any("A1" in str(m) for m in ws.merged_cells.ranges)
        assert merged, "Header cells were not merged"
        # Check formula
        assert ws["D4"].value == "=B4*C4", "Formula injection failed"
        # Check width
        assert ws.column_dimensions["A"].width == 20.0, "Column width not set"

    def test_multi_sheet_workbook_preservation_and_targeted_edits(self, client, tmp_path):
        """
        Scenario: An agent is tasked with updating a specific sheet ('Q4_Data') in a 
        multi-sheet corporate workbook without corrupting or altering other sheets ('Q1_Data', 'Dashboard').
        """
        input_path = tmp_path / "Corporate_Master.xlsx"
        output_path = tmp_path / "Corporate_Master_Updated.xlsx"
        
        # Create a multi-sheet workbook
        wb = openpyxl.Workbook()
        ws_q1 = wb.active
        ws_q1.title = "Q1_Data"
        ws_q1["A1"] = "Q1 Revenue"
        ws_q1["A2"] = 10000
        
        ws_dash = wb.create_sheet("Dashboard")
        ws_dash["A1"] = "Annual Overview"
        
        ws_q4 = wb.create_sheet("Q4_Data")
        ws_q4["A1"] = "Pending..."
        wb.save(input_path)
        
        payload = {
            "input_path": str(input_path),
            "output_path": str(output_path),
            "changes": [
                {"op": "set_value", "cell": "Q4_Data!A1", "value": "Q4 Revenue"},
                {"op": "set_value", "cell": "Q4_Data!A2", "value": 45000},
                {"op": "format_range", "range": "Q4_Data!A1:A1", "bold": True}
            ],
            "office_context": {"activeSheet": "Q4_Data", "preserveTemplate": True}
        }
        
        response = client.post("/skills/excel", json=payload)
        assert response.status_code == 200, f"Bridge failed: {response.text}"
        
        # Verify
        out_wb = openpyxl.load_workbook(output_path)
        assert "Q1_Data" in out_wb.sheetnames, "Template sheet Q1_Data was lost"
        assert "Dashboard" in out_wb.sheetnames, "Template sheet Dashboard was lost"
        assert out_wb["Q1_Data"]["A2"].value == 10000, "Untargeted sheet was modified"
        
        # Verify targeted edit
        assert out_wb["Q4_Data"]["A2"].value == 45000, "Targeted edit failed"
        assert out_wb["Q4_Data"]["A1"].font.bold is True, "Targeted formatting failed"


class TestCorporateDocumentDrafting:
    """Real-world scenarios for the Word Expert."""

    def test_compliant_contract_generation(self, client, tmp_path):
        """
        Scenario: An agent drafts a legal document. It must adhere to the corporate
        glossary (terminology enforcement) and replace a specific template section.
        """
        input_path = tmp_path / "Contract_Template.docx"
        output_path = tmp_path / "Final_Contract.docx"

        # Create a mock template
        doc = Document()
        doc.add_paragraph("{{LEGAL_TERMS}}", style="Normal")
        doc.save(input_path)

        payload = {
            "input_path": str(input_path),
            "output_path": str(output_path),
            "edits": [
                {
                    "op": "replace_section", 
                    "sectionId": "{{LEGAL_TERMS}}", 
                    # Simulating TS Invoker having applied Glossary correction
                    "text": "The Nexus Cloud Platform will be decommissioned.", 
                    "style": "Normal",
                    "metadata": {"glossaryCorrections": ["'Legacy System' -> 'Nexus Cloud Platform'"]}
                }
            ],
            "office_context": {
                "glossary": {"Legacy System": "Nexus Cloud Platform"},
                "preserveTemplate": True
            }
        }

        response = client.post("/skills/word", json=payload)
        assert response.status_code == 200, f"Bridge failed: {response.text}"
        
        data = response.json()
        assert data["template_preserved"] is True
        
        # Verify Metadata Transparency
        applied = data.get("applied_operations", [])
        assert len(applied) == 1
        assert "glossaryCorrections" in applied[0].get("metadata", {})

        # Verify Document Content
        final_doc = Document(output_path)
        content = "\n".join(p.text for p in final_doc.paragraphs)
        assert "Nexus Cloud Platform" in content, "Section replacement failed"
        assert "{{LEGAL_TERMS}}" not in content, "Placeholder was not removed"

    def test_executive_briefing_with_complex_tables_and_styled_spans(self, client, tmp_path):
        """
        Scenario: An agent creates an executive briefing containing a data table
        and performs targeted font styling (simulating markdown bolding on specific words).
        """
        output_path = tmp_path / "Briefing.docx"
        
        payload = {
            "output_path": str(output_path),
            "edits": [
                {"op": "insert_heading", "text": "Q4 Performance Overview", "level": 1},
                {"op": "insert_paragraph", "text": "The quarter exceeded expectations significantly in the EMEA region.", "style": "Normal"},
                {"op": "set_font", "target": "exceeded expectations", "bold": True, "size_pt": 12},
                {
                    "op": "insert_table", 
                    "rows": 3, 
                    "cols": 2, 
                    "data": [
                        ["Region", "Growth"],
                        ["EMEA", "25%"],
                        ["APAC", "18%"]
                    ],
                    "style": "Table Grid"
                }
            ]
        }
        
        response = client.post("/skills/word", json=payload)
        assert response.status_code == 200, f"Bridge failed: {response.text}"
        
        doc = Document(output_path)
        # Check heading
        assert doc.paragraphs[0].text == "Q4 Performance Overview"
        # Check targeted bolding
        target_para = doc.paragraphs[1]
        assert "exceeded expectations" in target_para.text
        bold_found = any(run.bold and "exceeded expectations" in run.text for run in target_para.runs)
        assert bold_found, "Targeted font styling failed"
        
        # Check table
        assert len(doc.tables) == 1
        table = doc.tables[0]
        assert table.cell(1, 0).text == "EMEA"
        assert table.cell(2, 1).text == "18%"

    def test_corporate_policy_update_with_strict_style_enforcement(self, client, tmp_path):
        """
        Scenario: An agent updates a legacy policy document. It searches for specific text,
        replaces it, and applies a strict corporate named style ('Heading 2').
        """
        input_path = tmp_path / "Policy_V1.docx"
        output_path = tmp_path / "Policy_V2.docx"
        
        # Create legacy doc
        doc = Document()
        doc.add_paragraph("This policy uses standard passwords.", style="Normal")
        doc.save(input_path)
        
        payload = {
            "input_path": str(input_path),
            "output_path": str(output_path),
            "edits": [
                {"op": "find_replace", "find": "standard passwords", "replace": "Multi-Factor Authentication (MFA)"},
                {"op": "apply_named_style", "target": "Multi-Factor Authentication (MFA)", "style": "Heading 2"}
            ],
            "office_context": {"preserveTemplate": True}
        }
        
        response = client.post("/skills/word", json=payload)
        assert response.status_code == 200, f"Bridge failed: {response.text}"
        
        final_doc = Document(output_path)
        content = "\n".join(p.text for p in final_doc.paragraphs)
        assert "Multi-Factor Authentication (MFA)" in content, "Find/Replace failed"
        assert "standard passwords" not in content, "Legacy text remained"
        
        # Verify style application
        target_para = [p for p in final_doc.paragraphs if "Multi-Factor Authentication (MFA)" in p.text][0]
        assert target_para.style.name == "Heading 2", "Named style was not applied"


class TestBrandCompliantPresentation:
    """Real-world scenarios for the PPT Expert."""

    def test_branded_pitch_deck_creation(self, client, tmp_path):
        """
        Scenario: An agent builds a pitch deck using advanced shapes (Diamond, Star)
        and strictly adheres to the corporate theme colors provided by the host context.
        """
        output_path = tmp_path / "Q3_Pitch_Deck.pptx"

        payload = {
            "output_path": str(output_path),
            "slides": [
                {"op": "add_title_slide", "title": "Q3 Strategy", "subtitle": "Omni-Office Zenith"},
                {"op": "add_slide", "title": "Milestones", "layout_index": 1},
                {
                    "op": "add_shape", 
                    "slide_index": 1, 
                    "shape_type": "diamond", 
                    "left": 100, "top": 100, "width": 150, "height": 150,
                    "fill_color": "primary"  # Theme color reference
                },
                {
                    "op": "add_shape", 
                    "slide_index": 1, 
                    "shape_type": "star", 
                    "left": 300, "top": 100, "width": 150, "height": 150,
                    "fill_color": "accent"  # Theme color reference
                }
            ],
            "office_context": {
                "themeColors": {"primary": "1E2761", "accent": "FF5733"}
            }
        }

        response = client.post("/skills/ppt", json=payload)
        assert response.status_code == 200, f"Bridge failed: {response.text}"
        
        data = response.json()
        assert data["slide_count"] == 2
        
        # Verify PowerPoint Structure
        prs = Presentation(output_path)
        assert len(prs.slides) == 2
        
        # Slide 0 is title, Slide 1 is Milestones + 2 Shapes (Title + Body + Diamond + Star = 4 shapes)
        slide_1_shapes = prs.slides[1].shapes
        assert len(slide_1_shapes) >= 3, "Shapes were not successfully injected"

    def test_sales_dashboard_with_media_and_presenter_notes(self, client, tmp_path):
        """
        Scenario: An agent generates a sales dashboard slide, changes the background color,
        and writes specific presenter notes for the speaker.
        """
        output_path = tmp_path / "Sales_Dashboard.pptx"
        image_path = tmp_path / "chart.png"
        
        # Create a dummy image
        from PIL import Image
        img = Image.new('RGB', (100, 100), color = 'red')
        img.save(image_path)
        
        payload = {
            "output_path": str(output_path),
            "slides": [
                {"op": "add_slide", "title": "Sales Funnel", "layout_index": 5},
                {"op": "set_background_color", "slide_index": 0, "hex_color": "secondary"},
                {"op": "add_image", "slide_index": 0, "image_path": str(image_path), "left_in": 2.0, "top_in": 2.0, "width_in": 4.0},
                {"op": "set_slide_notes", "slide_index": 0, "notes": "Remember to emphasize the 20% YoY growth."}
            ],
            "office_context": {
                "themeColors": {"secondary": "F3F4F6"}
            }
        }
        
        response = client.post("/skills/ppt", json=payload)
        assert response.status_code == 200, f"Bridge failed: {response.text}"
        
        prs = Presentation(output_path)
        assert len(prs.slides) == 1
        slide = prs.slides[0]
        
        # Check notes
        assert slide.has_notes_slide
        assert "20% YoY growth" in slide.notes_slide.notes_text_frame.text, "Presenter notes failed"
        
        # Check image (should be a picture shape, shape_type 13)
        pictures = [shape for shape in slide.shapes if getattr(shape, "shape_type", None) == 13]
        assert len(pictures) >= 1, "Image was not successfully inserted"

    def test_executive_summary_update_via_named_placeholders(self, client, tmp_path):
        """
        Scenario: An agent updates an existing Executive Summary slide by injecting
        specific text into named shapes (e.g., 'SummaryBox', 'KPIBox'), demonstrating
        precise contextual targeting.
        """
        input_path = tmp_path / "Exec_Summary_Draft.pptx"
        output_path = tmp_path / "Exec_Summary_Final.pptx"
        
        # Create draft presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
        
        # Add named shapes
        shape1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        shape1.name = "SummaryBox"
        shape1.text_frame.text = "[Insert Summary]"
        
        shape2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
        shape2.name = "KPIBox"
        shape2.text_frame.text = "[Insert KPI]"
        
        prs.save(input_path)
        
        payload = {
            "input_path": str(input_path),
            "output_path": str(output_path),
            "slides": [
                {"op": "insert_text", "slide_index": 0, "shape_name": "SummaryBox", "text": "Q4 performance was outstanding."},
                {"op": "insert_text", "slide_index": 0, "shape_name": "KPIBox", "text": "Revenue +35% YoY"},
                {"op": "set_font", "slide_index": 0, "shape_name": "KPIBox", "bold": True, "size_pt": 24}
            ],
            "office_context": {"preserveTemplate": True}
        }
        
        response = client.post("/skills/ppt", json=payload)
        assert response.status_code == 200, f"Bridge failed: {response.text}"
        
        final_prs = Presentation(output_path)
        final_slide = final_prs.slides[0]
        
        summary_shape = [s for s in final_slide.shapes if s.name == "SummaryBox"][0]
        kpi_shape = [s for s in final_slide.shapes if s.name == "KPIBox"][0]
        
        assert summary_shape.text_frame.text == "Q4 performance was outstanding."
        assert kpi_shape.text_frame.text == "Revenue +35% YoY"
        
        # Verify font formatting
        first_run = kpi_shape.text_frame.paragraphs[0].runs[0]
        assert first_run.font.bold is True, "Targeted font bolding failed"
        # font.size is in Pt, so it might be stored as Emu. We check if it's set.
        assert first_run.font.size is not None, "Targeted font sizing failed"
