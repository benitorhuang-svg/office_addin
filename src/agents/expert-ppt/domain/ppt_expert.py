from __future__ import annotations

import json
import os
import sys
from typing import Any, TypedDict, List, Optional, Dict

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from bridge_protocol import BaseExpertSkill

# ── Type Definitions ─────────────────────────────────────────────────────

class PPTOperation(TypedDict, total=False):
    op: str
    slide_index: Optional[int]
    title: Optional[str]
    subtitle: Optional[str]
    body: Optional[str]
    shape_type: Optional[str]
    text: Optional[str]
    left: Optional[float]
    top: Optional[float]
    width: Optional[float]
    height: Optional[float]
    fill_color: Optional[str]
    font_size_pt: Optional[int]
    shape_name: Optional[str]
    image_path: Optional[str]
    hex_color: Optional[str]
    notes: Optional[str]
    metadata: Optional[Dict[str, Any]]

class PPTPayload(TypedDict):
    input: Optional[str]
    input_path: Optional[str]
    output: str
    output_path: Optional[str]
    slides: List[PPTOperation]
    office_context: Optional[Dict[str, Any]]

# ── Brand defaults ────────────────────────────────────────────────────────
NEXUS_BLUE = RGBColor(0x00, 0x8C, 0xA1)
DEFAULT_FONT = "Calibri"
SUPPORTED_PPTX_SUFFIXES = {".pptx"}

class PPTMasterSkill(BaseExpertSkill):
    """
    Industrial PPT Design Master: High-fidelity layout and content orchestration.
    """
    def __init__(self, ppt_path: Optional[str] = None, office_context: Optional[Dict[str, Any]] = None) -> None:
        self.input_path = ppt_path or ""
        self.office_context = office_context or {}
        self.prs = self._load_presentation(ppt_path) if ppt_path else Presentation()

    def _op_add_slide(self, op: PPTOperation) -> None:
        layout_idx = op.get("layout_index", 1)
        layout_idx = min(layout_idx, len(self.prs.slide_layouts) - 1)
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)
        if slide.shapes.title and op.get("title"):
            slide.shapes.title.text = str(op["title"])
        if op.get("body") and len(slide.placeholders) > 1:
            body_ph = slide.placeholders[1]
            body_ph.text = str(op["body"])

    def _op_add_shape(self, op: PPTOperation) -> None:
        slide_idx = op.get("slide_index", 0)
        slide = self.prs.slides[slide_idx]
        left, top = Pt(op.get("left", 72)), Pt(op.get("top", 72))
        width, height = Pt(op.get("width", 240)), Pt(op.get("height", 120))
        shape_type = str(op.get("shape_type", "rectangle")).lower()

        if shape_type == "textbox":
            shape = slide.shapes.add_textbox(left, top, width, height)
        else:
            shape_enum = {
                "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                "oval": MSO_AUTO_SHAPE_TYPE.OVAL,
                "diamond": MSO_AUTO_SHAPE_TYPE.DIAMOND,
                "arrow": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                "star": MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
                "cloud": MSO_AUTO_SHAPE_TYPE.CLOUD,
            }.get(shape_type, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
            shape = slide.shapes.add_shape(shape_enum, left, top, width, height)
            fill_color = self._resolve_theme_color(op.get("fill_color"))
            if fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)

        if getattr(shape, "has_text_frame", False) and op.get("text"):
            shape.text_frame.text = str(op["text"])

    def _op_insert_text(self, op: PPTOperation) -> None:
        slide = self.prs.slides[op.get("slide_index", 0)]
        for shape in slide.shapes:
            if shape.name == op.get("shape_name") and shape.has_text_frame:
                shape.text_frame.text = str(op["text"])
                return

    def _op_set_font(self, op: PPTOperation) -> None:
        slide = self.prs.slides[op.get("slide_index", 0)]
        for shape in slide.shapes:
            if shape.name == op.get("shape_name") and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if op.get("size_pt"): run.font.size = Pt(op["size_pt"])
                        if op.get("bold") is not None: run.font.bold = op["bold"]

    def _op_add_image(self, op: PPTOperation) -> None:
        slide = self.prs.slides[op.get("slide_index", 0)]
        slide.shapes.add_picture(
            str(op["image_path"]), 
            Pt(op["left"]) if "left" in op else Inches(op.get("left_in", 1.0)),
            Pt(op["top"]) if "top" in op else Inches(op.get("top_in", 1.0))
        )

    def _op_add_title_slide(self, op: PPTOperation) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        if slide.shapes.title: slide.shapes.title.text = op.get("title", "")
        if len(slide.placeholders) > 1: slide.placeholders[1].text = op.get("subtitle", "")

    def _op_set_slide_notes(self, op: PPTOperation) -> None:
        self.prs.slides[op.get("slide_index", 0)].notes_slide.notes_text_frame.text = op.get("notes", "")

    def _op_set_background_color(self, op: PPTOperation) -> None:
        slide = self.prs.slides[op.get("slide_index", 0)]
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(self._resolve_theme_color(op.get("hex_color")) or "FFFFFF")

    _DISPATCH = {
        "add_slide": "_op_add_slide",
        "add_title_slide": "_op_add_title_slide",
        "add_shape": "_op_add_shape",
        "insert_text": "_op_insert_text",
        "set_slide_notes": "_op_set_slide_notes",
        "set_background_color": "_op_set_background_color",
        "set_font": "_op_set_font",
        "add_image": "_op_add_image",
    }

    def save_presentation(self, output_path: str) -> Dict[str, Any]:
        output_extension = self._get_extension(output_path)
        if output_extension not in SUPPORTED_PPTX_SUFFIXES:
            raise ValueError(f"Unsupported PowerPoint output format: {output_extension or '(missing extension)'}")
        parent_dir = os.path.dirname(output_path)
        if parent_dir and not os.path.exists(parent_dir):
            os.makedirs(parent_dir, exist_ok=True)
        self.prs.save(output_path)
        return {
            "status": "success", 
            "file": output_path, 
            "slide_count": len(self.prs.slides),
            "output_format": output_extension.lstrip("."),
            "office_context_received": bool(self.office_context),
            "template_preserved": bool(self.input_path),
        }

    def _load_presentation(self, ppt_path: str) -> Presentation:
        input_extension = self._get_extension(ppt_path)
        if input_extension == ".ppt":
            raise ValueError("Legacy .ppt files must be converted to .pptx before editing")
        if input_extension not in SUPPORTED_PPTX_SUFFIXES:
            raise ValueError(f"Unsupported PowerPoint input format: {input_extension or '(missing extension)'}")
        return Presentation(ppt_path)

    def _get_extension(self, file_path: str | None) -> str:
        return os.path.splitext(file_path or "")[1].lower()

    def _resolve_theme_color(self, value: Any) -> Optional[str]:
        if not value: return None
        theme_colors = self.office_context.get("themeColors", {})
        return theme_colors.get(value, str(value)).lstrip("#")

def run(payload: PPTPayload) -> Dict[str, Any]:
    input_path = payload.get("input") or payload.get("input_path")
    output_path = payload.get("output") or payload.get("output_path", "output.pptx")
    master = PPTMasterSkill(input_path, payload.get("office_context"))
    applied = master.apply_ops(payload.get("slides", []))
    for op in applied:
        if op.get("status") == "error":
            raise ValueError(op.get("message"))
    result = master.save_presentation(output_path)
    result["applied_operations"] = applied
    return result

if __name__ == "__main__":
    try:
        data = json.load(sys.stdin)
        print(json.dumps(run(data)))
    except Exception as e:
        print(json.dumps({"error": str(e)}), file=sys.stderr)
        sys.exit(1)
